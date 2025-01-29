from flask import Flask, render_template, request, redirect, url_for, send_from_directory
import os
import re
from docx2txt import process as doc_process
from pptx import Presentation
import pythoncom
import comtypes.client
from io import BytesIO
import base64

app = Flask(__name__)

# Configurations
app.config['UPLOAD_FOLDER'] = os.path.abspath('uploads')
app.config['GENERATED_FOLDER'] = os.path.abspath('generated')
app.config['THEMES_FOLDER'] = os.path.abspath('themes')  # Folder to store theme files

# Ensure necessary folders exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['GENERATED_FOLDER'], exist_ok=True)
os.makedirs(app.config['THEMES_FOLDER'], exist_ok=True)

@app.route('/home', methods=['GET', 'POST'])
def index():
    # Fetch available themes from the themes folder
    themes = [f for f in os.listdir(app.config['THEMES_FOLDER']) if f.endswith('.pptx')]

    if request.method == 'POST':
        # File upload handling
        file = request.files.get('docx_file')
        if not file or file.filename == '':
            return redirect(url_for('index'))
        
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filepath)

        # Theme selection handling
        selected_theme = request.form.get('theme')
        if not selected_theme or selected_theme not in themes:
            return redirect(url_for('index'))

        theme_path = os.path.join(app.config['THEMES_FOLDER'], selected_theme)

        # Process DOCX and generate PPT
        text = process_docx(filepath)
        cleaned_text = clean_unwanted_patterns(text)

        section_markers = ["Course", "Sem/AY", "Module No.", "Lesson Title", "Description of the Lesson",
                           "Intended Learning Outcomes", "Targets/ Objectives", "Learning Guide Questions",
                           "Lecture Guide", "Performance Task", "Learning Resources"]
        sections = extract_sections_and_content(cleaned_text, section_markers)

        ppt_file_path = os.path.join(app.config['GENERATED_FOLDER'], 'Document_Lesson_Themed.pptx')
        create_ppt_with_theme(sections, theme_path, ppt_file_path)

        pythoncom.CoInitialize()
        png_data_list = ppt_to_png(ppt_file_path)
        return render_template('index.html', themes=themes, png_data_list=png_data_list)

    return render_template('index.html', themes=themes)

@app.route('/')
def home():
    return render_template('home.html')

@app.route('/download/<filename>', methods=['POST'])
def download_ppt(filename):
    # Get the custom filename from the form
    custom_filename = request.form.get('custom_filename', 'Default_Presentation') + ".pptx"

    # Serve the file with the custom filename
    return send_from_directory(
        app.config['GENERATED_FOLDER'],
        filename,
        as_attachment=True,
        download_name=custom_filename
    )

def process_docx(docx_path):
    return doc_process(docx_path)

def clean_unwanted_patterns(text, additional_patterns=None):
    unwanted_patterns = [
        r"Date\s*\n?.*\d{4}",
        r"Week\s*Duration\s*\n?.*\d+",
        r"Online\s*Activities\s*\(Synchronous/.*",
        r"Offline\s*Activities.*",
        r"Face to Face Activities \(Synchronous/.*",
        r"e-Learning/Self-Paced.*",
        r"Online\s*Discussion\s*via\s*Google\s*Meet.*",
        r"Student Learning Strategies.*",
        r"Asynchronous.*",
        r"LSPU SELF-PACED LEARNING MODULE: TECHNOLOGY FOR TEACHING AND LEARNING.*",
        r"The online discussion.*",
        r"Performance Task Date.*",
        r"You will be directed.*",
        r"The online discussion will happen.*"
    ]
    if additional_patterns:
        unwanted_patterns.extend(additional_patterns)
    
    for pattern in unwanted_patterns:
        text = re.sub(pattern, "", text, flags=re.MULTILINE | re.IGNORECASE)
    
    text = re.sub(r'\n\s*\n', '\n\n', text).strip()
    return text

def extract_sections_and_content(text, section_markers):
    sections = {}
    current_section = None
    buffer = []
    
    for line in text.splitlines():
        line = line.strip()
        if any(line.startswith(marker) for marker in section_markers):
            if current_section:
                sections[current_section] = "\n".join(buffer)
            current_section = line
            buffer = []
        elif current_section:
            buffer.append(line)
    
    if current_section:
        sections[current_section] = "\n".join(buffer)
    
    return sections

def create_ppt_with_theme(sections, theme_file, ppt_file_path, word_limit=50):
    presentation = Presentation(theme_file)
    
    # Add title slide
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = sections.get("Lesson Title", "Untitled Document")

    # Add content slides
    for section, content in sections.items():
        if section != "Lesson Title":
            content_chunks = chunk_text(content, word_limit)
            for i, chunk in enumerate(content_chunks):
                slide_layout = presentation.slide_layouts[1]
                slide = presentation.slides.add_slide(slide_layout)
                slide.shapes.title.text = f"{section} (Part {i+1})" if i > 0 else section
                slide.placeholders[1].text = chunk

    presentation.save(ppt_file_path)

def chunk_text(text, word_limit):
    words = text.split()
    chunks, chunk = [], []
    for word in words:
        chunk.append(word)
        if len(chunk) >= word_limit:
            chunks.append(" ".join(chunk))
            chunk = []
    if chunk:
        chunks.append(" ".join(chunk))
    return chunks

def ppt_to_png(input_ppt):
    if not os.path.exists(input_ppt):
        print(f"Error: The file {input_ppt} does not exist.")
        return []

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(input_ppt)

    png_data_list = []
    for slide_index in range(1, presentation.Slides.Count + 1):
        slide = presentation.Slides(slide_index)
        temp_file = BytesIO()
        temp_path = os.path.abspath(f"temp_slide_{slide_index}.png")
        slide.Export(temp_path, "PNG")
        
        with open(temp_path, "rb") as temp_img_file:
            temp_file.write(temp_img_file.read())
        
        temp_file.seek(0)
        png_data = base64.b64encode(temp_file.getvalue()).decode('utf-8')
        png_data_list.append(png_data)
        
        os.remove(temp_path)

    presentation.Close()
    powerpoint.Quit()

    return png_data_list

if __name__ == '__main__':
    app.run(debug=True)
